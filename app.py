"""
INT Marketing Translation Agent — Streamlit POC
================================================
A proof-of-concept for the automated promotional material localization pipeline.

Run with:
    streamlit run app.py

Requirements:
    pip install streamlit anthropic python-docx PyPDF2 python-pptx
"""

import streamlit as st
import anthropic
import boto3
import os
import json
import time
import re
import zipfile
import io
from pathlib import Path
from datetime import datetime
from typing import Optional, Annotated, TypedDict
import operator

# ── LangGraph availability ───────────────────────────────────────────────────
try:
    from langgraph.graph import StateGraph, END, START
    LANGGRAPH_AVAILABLE = True
except ImportError:
    LANGGRAPH_AVAILABLE = False
    StateGraph = END = START = None

# ── Page config ─────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="INT Mkt Translation Agent",
    page_icon="🌐",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Custom CSS ───────────────────────────────────────────────────────────────
st.markdown("""
<style>
  /* Global */
  [data-testid="stAppViewContainer"] { background: #F7F8FC; }
  [data-testid="stSidebar"] { background: #0F2240; }
  [data-testid="stSidebar"] * { color: #E8EFF8 !important; }
  [data-testid="stSidebar"] .stSelectbox label,
  [data-testid="stSidebar"] .stMultiselect label,
  [data-testid="stSidebar"] .stTextInput label { color: #93B4D4 !important; font-size: 12px !important; text-transform: uppercase; letter-spacing: 0.08em; }
  [data-testid="stSidebar"] h2 { color: #FFFFFF !important; font-size: 18px !important; margin-bottom: 0.25rem !important; }
  [data-testid="stSidebar"] p { color: #93B4D4 !important; font-size: 12px !important; }

  /* Logo area */
  .logo-block { padding: 1rem 1rem 0.5rem; border-bottom: 1px solid rgba(255,255,255,0.1); margin-bottom: 1rem; }
  .logo-title { font-size: 16px; font-weight: 700; color: #FFFFFF !important; letter-spacing: 0.02em; }
  .logo-sub { font-size: 11px; color: #93B4D4 !important; margin-top: 2px; }

  /* Main header */
  .main-header { background: linear-gradient(135deg, #0F2240 0%, #185FA5 100%);
    border-radius: 12px; padding: 1.5rem 2rem; margin-bottom: 1.5rem; color: white; }
  .main-header h1 { color: white !important; font-size: 26px !important; margin: 0 !important; }
  .main-header p { color: rgba(255,255,255,0.75) !important; margin: 4px 0 0 !important; font-size: 14px !important; }

  /* Metric cards */
  .metric-row { display: flex; gap: 12px; margin-bottom: 1.25rem; }
  .metric-card { flex: 1; background: white; border-radius: 10px; padding: 1rem 1.25rem;
    border: 0.5px solid #E0E8F4; box-shadow: 0 1px 4px rgba(0,0,0,0.04); }
  .metric-card .mc-label { font-size: 11px; color: #7B8FA6; text-transform: uppercase; letter-spacing: 0.08em; margin-bottom: 4px; }
  .metric-card .mc-val { font-size: 28px; font-weight: 700; color: #0F2240; line-height: 1; }
  .metric-card .mc-sub { font-size: 12px; color: #7B8FA6; margin-top: 3px; }

  /* File cards */
  .file-card { background: white; border: 0.5px solid #E0E8F4; border-radius: 10px;
    padding: 1rem 1.25rem; margin-bottom: 0.6rem; cursor: pointer;
    transition: border-color 0.15s, box-shadow 0.15s; }
  .file-card:hover { border-color: #185FA5; box-shadow: 0 2px 8px rgba(24,95,165,0.1); }
  .file-card.selected { border-color: #185FA5; border-width: 1.5px; background: #F0F6FF; }
  .file-card .fc-name { font-size: 14px; font-weight: 600; color: #0F2240; margin-bottom: 2px; }
  .file-card .fc-meta { font-size: 12px; color: #7B8FA6; }
  .file-card .fc-type { display: inline-block; font-size: 10px; font-weight: 600;
    padding: 2px 7px; border-radius: 4px; margin-right: 6px; text-transform: uppercase; letter-spacing: 0.06em; }
  .fc-type-txt  { background: #E8F4E8; color: #2D7A2D; }
  .fc-type-docx { background: #E8F0FF; color: #2D50A0; }
  .fc-type-pdf  { background: #FFE8E8; color: #A02D2D; }
  .fc-type-pptx { background: #FFF0E0; color: #A06020; }

  /* Locale flag chips */
  .locale-chip { display: inline-flex; align-items: center; gap: 5px;
    background: white; border: 0.5px solid #D0DCF0; border-radius: 20px;
    padding: 4px 12px; margin: 3px; font-size: 13px; }
  .locale-chip.active { background: #185FA5; border-color: #185FA5; color: white; }
  .locale-chip .flag { font-size: 16px; }

  /* Progress / status */
  .status-pill { display: inline-block; font-size: 11px; font-weight: 600;
    padding: 3px 10px; border-radius: 20px; }
  .status-pending  { background: #FFF3CD; color: #856404; }
  .status-running  { background: #CCE5FF; color: #004085; }
  .status-done     { background: #D4EDDA; color: #155724; }
  .status-flagged  { background: #F8D7DA; color: #721C24; }

  /* Translation result card */
  .result-card { background: white; border: 0.5px solid #E0E8F4; border-radius: 10px;
    padding: 1.25rem; margin-bottom: 0.75rem; }
  .result-card .rc-header { display: flex; align-items: center; justify-content: space-between; margin-bottom: 0.75rem; }
  .result-card .rc-locale { font-size: 15px; font-weight: 700; color: #0F2240; }
  .result-card .rc-content { font-size: 13px; color: #2C3E50; line-height: 1.7;
    background: #F7F9FC; border-radius: 8px; padding: 1rem; max-height: 320px; overflow-y: auto; }
  .result-card .rc-flags { margin-top: 0.75rem; }
  .flag-item { background: #FFF3CD; border: 0.5px solid #F0C040; border-radius: 6px;
    padding: 0.5rem 0.75rem; margin-bottom: 0.4rem; font-size: 12px; }
  .flag-item .fi-reason { font-weight: 700; color: #856404; margin-right: 6px; }

  /* Section titles */
  .section-title { font-size: 13px; font-weight: 700; color: #7B8FA6;
    text-transform: uppercase; letter-spacing: 0.1em; margin: 1.25rem 0 0.6rem; }

  /* Audit log */
  .audit-row { font-size: 12px; color: #4A5568; padding: 0.4rem 0;
    border-bottom: 0.5px solid #EEF2F8; }
  .audit-row .ar-time { color: #7B8FA6; margin-right: 8px; font-family: monospace; }
  .audit-row .ar-action { color: #185FA5; font-weight: 600; margin-right: 6px; }

  /* Folder tree */
  .folder-node { font-size: 13px; padding: 3px 0; }
  .folder-node .fn-icon { margin-right: 6px; }
  .folder-node.is-dir { font-weight: 600; color: #0F2240; }
  .folder-node.is-file { color: #4A5568; padding-left: 1.5rem; }
  .folder-node.is-translated { color: #0F6E56; }

  /* Confidence bar */
  .conf-bar-wrap { height: 6px; background: #E8EFF4; border-radius: 3px; margin-top: 4px; }
  .conf-bar { height: 100%; border-radius: 3px; }

  /* Tab overrides */
  .stTabs [data-baseweb="tab-list"] { gap: 4px; background: transparent; border-bottom: 1px solid #E0E8F4; }
  .stTabs [data-baseweb="tab"] { background: transparent; border-radius: 8px 8px 0 0;
    font-size: 13px; font-weight: 500; padding: 0.5rem 1rem; color: #7B8FA6; }
  .stTabs [aria-selected="true"] { background: white !important; color: #185FA5 !important;
    border: 0.5px solid #E0E8F4 !important; border-bottom: none !important; }

  /* Spinner override */
  [data-testid="stSpinner"] { color: #185FA5; }

  /* Expander */
  details summary { font-size: 13px; font-weight: 600; color: #185FA5; }
</style>
""", unsafe_allow_html=True)

# ── Constants ─────────────────────────────────────────────────────────────────
LANGUAGES = {
    "fr-FR": {"name": "French",     "country": "France",         "flag": "🇫🇷"},
    "de-DE": {"name": "German",     "country": "Germany",        "flag": "🇩🇪"},
    "pt-BR": {"name": "Portuguese", "country": "Brazil",         "flag": "🇧🇷"},
    "ja-JP": {"name": "Japanese",   "country": "Japan",          "flag": "🇯🇵"},
    "es-ES": {"name": "Spanish",    "country": "Spain",          "flag": "🇪🇸"},
    "it-IT": {"name": "Italian",    "country": "Italy",          "flag": "🇮🇹"},
    "nl-NL": {"name": "Dutch",      "country": "Netherlands",    "flag": "🇳🇱"},
    "ko-KR": {"name": "Korean",     "country": "South Korea",    "flag": "🇰🇷"},
    "zh-CN": {"name": "Chinese",    "country": "China",          "flag": "🇨🇳"},
    "pl-PL": {"name": "Polish",     "country": "Poland",         "flag": "🇵🇱"},
    "ar-SA": {"name": "Arabic",     "country": "Saudi Arabia",   "flag": "🇸🇦"},
    "ru-RU": {"name": "Russian",    "country": "Russia",         "flag": "🇷🇺"},
}

BASE_DIR = Path(__file__).parent
SOURCE_DIR = BASE_DIR / "source_materials"
TRANSLATED_DIR = BASE_DIR / "translated"

# AWS Bedrock Configuration
AWS_PROFILE_NAME = "foresight"
AWS_REGION = "us-east-1"
AWS_MODEL_ID = "us.anthropic.claude-sonnet-4-20250514-v1:0"

COSTAR_SYSTEM_PROMPT = """You are the INT Marketing Translation Agent for a global pharmaceutical company.

## CONTEXT
You process approved promotional materials (campaign decks, product brochures, regulatory communications) from the INT Marketing team and translate them into target country languages while preserving all content integrity.

## OBJECTIVE
Translate the provided document chunk into the target locale. Return ONLY a valid JSON object with no preamble or markdown fencing.

## STYLE
- Marketing copy: natural, fluent, authored in the target language — not a word-for-word translation
- Regulatory language (contraindications, safety statements, clinical claims): translate faithfully, do not rephrase
- Brand/product names (CARDIOMAX, cardiomaxib, PharmaCorp): preserve exactly as written
- Clinical data (trial names, statistics, p-values, confidence intervals): preserve exactly

## TONE
Translated documents should feel authored in the target language. The reader should not detect the source language.

## FLAGGING RULES
Any segment containing:
- Regulatory claims or contraindications → flag as "regulatory_claim"
- Clinical trial data or statistics → flag as "clinical_data"
- Safety statements → flag as "regulatory_claim"
- Text you are less than 85% confident about → flag as "low_confidence"
- Brand terms with no clear local equivalent → flag as "brand_term_missing"

## RESPONSE FORMAT
Return ONLY this JSON schema (no markdown, no preamble):
{
  "translated_content": "full translated text",
  "source_locale": "en-US",
  "target_locale": "<locale>",
  "confidence_score": <0.0-1.0>,
  "flagged_segments": [
    {
      "segment_id": "<id>",
      "reason": "regulatory_claim|clinical_data|low_confidence|brand_term_missing",
      "original": "<source text>",
      "translation": "<translated text>",
      "confidence": <0.0-1.0>
    }
  ],
  "glossary_terms_applied": ["list of brand/product terms preserved"],
  "translation_notes": "brief reviewer notes"
}"""

# ── Agent System Prompt (tool-use mode) ──────────────────────────────────────
AGENT_SYSTEM_PROMPT = """You are the INT Marketing Translation Agent for a global pharmaceutical company.

## CONTEXT
You process approved promotional materials from the INT Marketing team and translate them into target country languages while preserving all content integrity.

## OBJECTIVE
Translate the provided document into the target locale using the available tools.

## STYLE
- Marketing copy: natural, fluent, authored in the target language — not a word-for-word translation
- Regulatory language (contraindications, safety statements, clinical claims): translate faithfully, do not rephrase
- Brand/product names (CARDIOMAX, cardiomaxib, PharmaCorp): preserve exactly as written
- Clinical data (trial names, statistics, p-values, confidence intervals): preserve exactly

## WORKFLOW
1. Translate the full document into the target language
2. Call flag_segment for each regulatory claim, clinical data reference, safety statement, or uncertain text
3. Call apply_glossary_term for each brand/product name you preserve
4. If images are provided, analyze them and incorporate any text they contain into the translation
5. Call submit_translation with your completed translation

## TONE
Translated documents should feel authored in the target language. The reader should not detect the source language."""

# ── Agent Tool Definitions ────────────────────────────────────────────────────
TRANSLATION_TOOLS = [
    {
        "name": "flag_segment",
        "description": "Flag a translated segment for human review due to regulatory, clinical, or quality concerns.",
        "input_schema": {
            "type": "object",
            "properties": {
                "segment_id": {"type": "string", "description": "Unique ID for this segment"},
                "reason": {
                    "type": "string",
                    "enum": ["regulatory_claim", "clinical_data", "low_confidence", "brand_term_missing"],
                    "description": "Reason for flagging"
                },
                "original":    {"type": "string", "description": "Original source text"},
                "translation": {"type": "string", "description": "Translated text"},
                "confidence":  {"type": "number", "description": "Confidence score 0.0-1.0"}
            },
            "required": ["segment_id", "reason", "original", "translation", "confidence"]
        }
    },
    {
        "name": "apply_glossary_term",
        "description": "Record a brand/product term that was preserved exactly as specified.",
        "input_schema": {
            "type": "object",
            "properties": {"term": {"type": "string", "description": "The brand/product term preserved"}},
            "required": ["term"]
        }
    },
    {
        "name": "submit_translation",
        "description": "Submit the final completed translation. Call this once when translation is done.",
        "input_schema": {
            "type": "object",
            "properties": {
                "translated_content": {"type": "string", "description": "Full translated document text"},
                "confidence_score":   {"type": "number", "description": "Overall confidence 0.0-1.0"},
                "translation_notes":  {"type": "string", "description": "Brief notes for the reviewer"}
            },
            "required": ["translated_content", "confidence_score", "translation_notes"]
        }
    }
]

# ── LangGraph Node System Prompts ─────────────────────────────────────────────
_TRANSLATOR_SYSTEM = """You are an expert pharmaceutical marketing translator.
Your role: produce a fluent, accurate translation into the target language.

Rules:
- Brand names (CARDIOMAX, cardiomaxib, PharmaCorp): preserve exactly as written
- Clinical data (statistics, p-values, trial names): preserve exactly
- Regulatory language: translate faithfully, do not rephrase or simplify
- Marketing copy: natural, fluent — feel authored in the target language, not translated

If QA feedback from a prior attempt is provided, address those specific issues.

Return ONLY valid JSON (no markdown fencing):
{"translated_content": "full translation", "glossary_terms": ["term1"], "confidence_score": 0.92}"""

_REGULATORY_SYSTEM = """You are a pharmaceutical regulatory compliance specialist.
Your role: review a translation and identify content requiring mandatory human review.

Flag segments containing:
- Regulatory claims or contraindications → "regulatory_claim"
- Clinical trial data or statistics → "clinical_data"
- Safety statements or warnings → "regulatory_claim"
- Text with confidence below 85% → "low_confidence"
- Brand terms with no clear local equivalent → "brand_term_missing"

Return ONLY valid JSON (no markdown fencing):
{"flagged_segments": [{"segment_id": "s1", "reason": "regulatory_claim", "original": "...", "translation": "...", "confidence": 0.9}], "regulatory_notes": "brief summary"}"""

_QA_SYSTEM = """You are a pharmaceutical translation quality assurance specialist.
Your role: final quality check before the translation is delivered.

Evaluate: completeness, accuracy, fluency, brand consistency, and regulatory adequacy.
Only request revision for genuine quality issues, not minor stylistic preferences.

Return ONLY valid JSON (no markdown fencing):
{"qa_feedback": "specific issues or 'Translation meets quality standards'", "revision_needed": false, "translation_notes": "summary for reviewer", "final_confidence_score": 0.92}"""


# ── LangGraph State Schema ─────────────────────────────────────────────────────
class TranslationState(TypedDict):
    source_content:    str
    images:            list
    locale:            str
    filename:          str
    lang_name:         str
    country:           str
    draft_translation: str
    flagged_segments:  Annotated[list, operator.add]
    glossary_terms:    Annotated[list, operator.add]
    regulatory_notes:  str
    qa_feedback:       str
    confidence_score:  float
    translation_notes: str
    revision_needed:   bool
    revision_count:    int
    errors:            Annotated[list, operator.add]


# ── LangGraph LLM Helper ──────────────────────────────────────────────────────
def _llm_call(client, system: str, user_content, use_bedrock: bool) -> str:
    """Single LLM call for a LangGraph node. Accepts str or list of content blocks."""
    if isinstance(user_content, str):
        user_content = [{"type": "text", "text": user_content}]
    if use_bedrock:
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 4000,
            "system": system,
            "messages": [{"role": "user", "content": user_content}]
        })
        resp = client.invoke_model(modelId=AWS_MODEL_ID, body=body)
        raw  = json.loads(resp["body"].read())["content"][0]["text"].strip()
    else:
        resp = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4000,
            system=system,
            messages=[{"role": "user", "content": user_content}]
        )
        raw = resp.content[0].text.strip()
    raw = re.sub(r'^```(?:json)?\s*', '', raw)
    raw = re.sub(r'\s*```$', '', raw)
    return raw


# ── LangGraph Pipeline ────────────────────────────────────────────────────────
def build_translation_graph(ai_client, use_bedrock: bool):
    """Compile the 3-node LangGraph translation pipeline."""
    if not LANGGRAPH_AVAILABLE:
        raise ImportError("langgraph not installed. Run: pip install langgraph>=0.2.0")

    def call(system: str, user_content) -> str:
        return _llm_call(ai_client, system, user_content, use_bedrock)

    # ── Node 1: Translator ────────────────────────────────────────────────────
    def translate_node(state: TranslationState) -> dict:
        qa_ctx = (
            f"\n\nQA FEEDBACK FROM PREVIOUS ATTEMPT — address these issues:\n{state['qa_feedback']}"
            if state.get("qa_feedback") and state.get("revision_count", 0) > 0 else ""
        )
        prompt = (
            f"Translate from English (en-US) to {state['lang_name']} ({state['locale']}) "
            f"for {state['country']}.\n\nDOCUMENT: {state['filename']}{qa_ctx}\n\n"
            f"SOURCE CONTENT:\n---\n{state['source_content'][:6000]}\n---"
        )
        blocks = [{"type": "text", "text": prompt}]
        for img in state.get("images", [])[:5]:
            blocks.append({"type": "image",
                            "source": {"type": "base64",
                                       "media_type": img["media_type"],
                                       "data": img["data"]}})
        try:
            parsed = json.loads(call(_TRANSLATOR_SYSTEM, blocks))
            return {
                "draft_translation": parsed.get("translated_content", ""),
                "glossary_terms":    parsed.get("glossary_terms", []),
                "confidence_score":  parsed.get("confidence_score", 0.0),
                "revision_count":    state.get("revision_count", 0) + 1,
            }
        except Exception as e:
            return {"errors": [f"translate_node: {e}"],
                    "revision_count": state.get("revision_count", 0) + 1}

    # ── Node 2: Regulatory Reviewer ───────────────────────────────────────────
    def regulatory_review_node(state: TranslationState) -> dict:
        if not state.get("draft_translation"):
            return {"regulatory_notes": "No translation to review.", "flagged_segments": []}
        prompt = (
            f"Review this translation ({state['locale']}) for regulatory/clinical content.\n\n"
            f"SOURCE:\n{state['source_content'][:3000]}\n\n"
            f"TRANSLATION:\n{state['draft_translation'][:3000]}"
        )
        try:
            parsed = json.loads(call(_REGULATORY_SYSTEM, prompt))
            return {
                "flagged_segments": parsed.get("flagged_segments", []),
                "regulatory_notes": parsed.get("regulatory_notes", ""),
            }
        except Exception as e:
            return {"errors": [f"regulatory_review_node: {e}"],
                    "regulatory_notes": f"Review error: {e}", "flagged_segments": []}

    # ── Node 3: QA Checker ────────────────────────────────────────────────────
    def qa_node(state: TranslationState) -> dict:
        if not state.get("draft_translation"):
            return {"revision_needed": False,
                    "translation_notes": "No translation produced.", "qa_feedback": ""}
        flags_json = json.dumps(state.get("flagged_segments", [])[:5], indent=2)
        prompt = (
            f"Quality-check this {state['locale']} translation.\n\n"
            f"REGULATORY NOTES: {state.get('regulatory_notes', 'None')}\n"
            f"FLAGGED SEGMENTS:\n{flags_json}\n\n"
            f"SOURCE (excerpt):\n{state['source_content'][:2000]}\n\n"
            f"TRANSLATION:\n{state['draft_translation'][:3000]}"
        )
        try:
            parsed = json.loads(call(_QA_SYSTEM, prompt))
            revision_needed = parsed.get("revision_needed", False)
            if state.get("revision_count", 0) >= 2:   # cap at 2 revision passes
                revision_needed = False
            return {
                "qa_feedback":       parsed.get("qa_feedback", ""),
                "revision_needed":   revision_needed,
                "translation_notes": parsed.get("translation_notes", ""),
                "confidence_score":  parsed.get("final_confidence_score",
                                                state.get("confidence_score", 0.0)),
            }
        except Exception as e:
            return {"errors": [f"qa_node: {e}"], "revision_needed": False,
                    "translation_notes": f"QA error: {e}", "qa_feedback": ""}

    # ── Conditional edge: QA → revise or done ─────────────────────────────────
    def should_revise(state: TranslationState) -> str:
        return "translate" if state.get("revision_needed") else END

    # ── Build graph ───────────────────────────────────────────────────────────
    graph = StateGraph(TranslationState)
    graph.add_node("translate",         translate_node)
    graph.add_node("regulatory_review", regulatory_review_node)
    graph.add_node("qa_check",          qa_node)
    graph.add_edge(START,               "translate")
    graph.add_edge("translate",         "regulatory_review")
    graph.add_edge("regulatory_review", "qa_check")
    graph.add_conditional_edges("qa_check", should_revise)
    return graph.compile()


def translate_document_langgraph_bedrock(bedrock_client, content: str, images: list,
                                          locale: str, filename: str) -> dict:
    """LangGraph 3-node pipeline: Translator → Regulatory → QA (via AWS Bedrock)."""
    lang_info = LANGUAGES.get(locale, {})
    graph     = build_translation_graph(bedrock_client, use_bedrock=True)
    initial   = TranslationState(
        source_content=content, images=images, locale=locale, filename=filename,
        lang_name=lang_info.get("name", locale), country=lang_info.get("country", ""),
        draft_translation="", flagged_segments=[], glossary_terms=[],
        regulatory_notes="", qa_feedback="", confidence_score=0.0,
        translation_notes="", revision_needed=False, revision_count=0, errors=[]
    )
    final = graph.invoke(initial)
    return {
        "translated_content":     final.get("draft_translation", ""),
        "source_locale":          "en-US",
        "target_locale":          locale,
        "confidence_score":       final.get("confidence_score", 0.0),
        "flagged_segments":       final.get("flagged_segments", []),
        "glossary_terms_applied": final.get("glossary_terms", []),
        "translation_notes":      final.get("translation_notes", ""),
    }


def translate_document_langgraph(api_client: anthropic.Anthropic, content: str, images: list,
                                   locale: str, filename: str) -> dict:
    """LangGraph 3-node pipeline: Translator → Regulatory → QA (via Anthropic API)."""
    lang_info = LANGUAGES.get(locale, {})
    graph     = build_translation_graph(api_client, use_bedrock=False)
    initial   = TranslationState(
        source_content=content, images=images, locale=locale, filename=filename,
        lang_name=lang_info.get("name", locale), country=lang_info.get("country", ""),
        draft_translation="", flagged_segments=[], glossary_terms=[],
        regulatory_notes="", qa_feedback="", confidence_score=0.0,
        translation_notes="", revision_needed=False, revision_count=0, errors=[]
    )
    final = graph.invoke(initial)
    return {
        "translated_content":     final.get("draft_translation", ""),
        "source_locale":          "en-US",
        "target_locale":          locale,
        "confidence_score":       final.get("confidence_score", 0.0),
        "flagged_segments":       final.get("flagged_segments", []),
        "glossary_terms_applied": final.get("glossary_terms", []),
        "translation_notes":      final.get("translation_notes", ""),
    }


# ── Session state init ────────────────────────────────────────────────────────
def init_state():
    defaults = {
        "api_key": "",
        "translation_results": {},
        "audit_log": [],
        "selected_file": None,
        "job_running": False,
        "uploaded_files_cache": {},
        "llm_provider": "AWS Bedrock",  # Default to AWS Bedrock
        "bedrock_client": None,
        "agent_mode": "LangGraph",
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

init_state()

# ── Helpers ───────────────────────────────────────────────────────────────────
def read_file_content(path: Path) -> str:
    """Extract text from .txt, .docx, .pdf, .pptx files."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".txt":
            return path.read_text(encoding="utf-8", errors="replace")
        elif suffix == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                return path.read_text(encoding="utf-8", errors="replace")
        elif suffix == ".pdf":
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(str(path))
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            except ImportError:
                return "[PDF reading requires: pip install PyPDF2]"
        elif suffix in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                return "\n".join(texts)
            except ImportError:
                return "[PPTX reading requires: pip install python-pptx]"
        else:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[Error reading file: {e}]"


def read_uploaded_content(uploaded_file) -> str:
    """Read content from a Streamlit UploadedFile object."""
    name = uploaded_file.name
    suffix = Path(name).suffix.lower()
    raw = uploaded_file.read()
    if suffix == ".txt":
        return raw.decode("utf-8", errors="replace")
    elif suffix == ".docx":
        try:
            from docx import Document
            import io
            doc = Document(io.BytesIO(raw))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return raw.decode("utf-8", errors="replace")
    elif suffix == ".pdf":
        try:
            import PyPDF2, io
            reader = PyPDF2.PdfReader(io.BytesIO(raw))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except ImportError:
            return "[PDF reading requires: pip install PyPDF2]"
    elif suffix in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(raw))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            return "[PPTX reading requires: pip install python-pptx]"
    else:
        return raw.decode("utf-8", errors="replace")


# ── Vision & Image Extraction ─────────────────────────────────────────────────
def extract_pptx_with_images(path: Path) -> tuple:
    """Extract text and base64 images from a PPTX file. Returns (text, [image_dicts])."""
    import base64
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(str(path))
        texts, images = [], []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        ext  = shape.image.ext.lower()
                        mime = {"png": "image/png", "jpg": "image/jpeg",
                                "jpeg": "image/jpeg", "gif": "image/gif",
                                "webp": "image/webp"}.get(ext, "image/png")
                        images.append({"slide": slide_num, "media_type": mime,
                                       "data": base64.standard_b64encode(blob).decode()})
                    except Exception:
                        pass
            if slide_texts:
                texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts), images
    except ImportError:
        return "[python-pptx not installed]", []
    except Exception as e:
        return f"[PPTX error: {e}]", []


def extract_docx_with_images(path: Path) -> tuple:
    """Extract text and base64 images from a DOCX file. Returns (text, [image_dicts])."""
    import base64
    try:
        from docx import Document
        doc = Document(str(path))
        text   = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    blob = rel.target_part.blob
                    mime = rel.target_part.content_type
                    if not mime.startswith("image/"):
                        mime = "image/png"
                    images.append({"media_type": mime,
                                   "data": base64.standard_b64encode(blob).decode()})
                except Exception:
                    pass
        return text, images
    except ImportError:
        return "[python-docx not installed]", []
    except Exception as e:
        return f"[DOCX error: {e}]", []


def read_file_with_images(path: Path) -> tuple:
    """Returns (text_content, images_list) for all supported file types."""
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return extract_pptx_with_images(path)
    elif suffix == ".docx":
        return extract_docx_with_images(path)
    else:
        return read_file_content(path), []


# ── Structured Segment Extraction (for native-format output) ─────────────────

def extract_pptx_segments(path: Path) -> list:
    """Extract per-paragraph text segments from a PPTX for structured translation.
    Handles regular text frames AND table cells.
    Returns list of {"id": "...", "text": "..."}.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(str(path))
        segments = []
        for si, slide in enumerate(prs.slides):
            for shi, shape in enumerate(slide.shapes):
                # ── Regular text frame ──────────────────────────────────────
                if shape.has_text_frame:
                    for pi, para in enumerate(shape.text_frame.paragraphs):
                        text = para.text.strip()
                        if text:
                            segments.append({"id": f"s{si}_sh{shi}_p{pi}", "text": text})
                # ── Table cells ─────────────────────────────────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    for ri, row in enumerate(shape.table.rows):
                        for ci, cell in enumerate(row.cells):
                            for pi, para in enumerate(cell.text_frame.paragraphs):
                                text = para.text.strip()
                                if text:
                                    segments.append({
                                        "id": f"s{si}_sh{shi}_r{ri}c{ci}_p{pi}",
                                        "text": text
                                    })
        return segments
    except Exception:
        return []


def extract_docx_segments(path: Path) -> list:
    """Extract per-paragraph text segments from a DOCX for structured translation.
    Handles both regular paragraphs AND table cells.
    Returns list of {"id": "...", "text": "..."}.
    """
    try:
        from docx import Document
        doc = Document(str(path))
        segments = []
        # Regular paragraphs
        for pi, p in enumerate(doc.paragraphs):
            if p.text.strip():
                segments.append({"id": f"p{pi}", "text": p.text.strip()})
        # Table cells
        for ti, table in enumerate(doc.tables):
            for ri, row in enumerate(table.rows):
                for ci, cell in enumerate(row.cells):
                    for pi, para in enumerate(cell.paragraphs):
                        if para.text.strip():
                            segments.append({
                                "id": f"t{ti}_r{ri}_c{ci}_p{pi}",
                                "text": para.text.strip()
                            })
        return segments
    except Exception:
        return []


def _build_segment_prompt(segments: list, locale: str, filename: str) -> str:
    lang_info = LANGUAGES.get(locale, {})
    input_json = json.dumps(
        [{"id": s["id"], "text": s["text"]} for s in segments],
        ensure_ascii=False
    )
    return (
        f"Translate the following JSON array of text segments from English (en-US) to "
        f"{lang_info.get('name', locale)} ({locale}) for {lang_info.get('country', '')}.\n\n"
        f"Document: {filename}\n\n"
        f"Rules:\n"
        f"- Marketing copy: natural and fluent in the target language\n"
        f"- Brand names (CARDIOMAX, cardiomaxib, PharmaCorp): preserve exactly as written\n"
        f"- Clinical data, statistics, p-values, trial names: preserve exactly\n"
        f"- Return ONLY a JSON array in the EXACT same order with \"text\" replaced by the translation\n"
        f"- Keep all \"id\" values unchanged. Do NOT wrap output in markdown fences.\n\n"
        f"Input:\n{input_json}"
    )


def _parse_segment_response(raw: str) -> dict:
    """Strip markdown fences and parse JSON array → {id: translated_text}.
    Robust: extracts the JSON array even if Claude wraps it in explanatory text.
    """
    raw = re.sub(r'^```[a-z]*\n?', '', raw.strip())
    raw = re.sub(r'\n?```$', '', raw.strip())
    # If Claude added preamble/postamble, extract the JSON array portion
    match = re.search(r'\[[\s\S]*\]', raw)
    if match:
        raw = match.group(0)
    return {item["id"]: item["text"] for item in json.loads(raw)}


_SEGMENT_CHUNK_SIZE = 60  # Translate at most this many segments per API call


def translate_segments_bedrock(bedrock_client, segments: list, locale: str, filename: str) -> dict:
    """Translate structured segments via AWS Bedrock. Chunks large presentations.
    Returns {id: translated_text}.
    """
    all_translations: dict = {}
    for chunk_start in range(0, len(segments), _SEGMENT_CHUNK_SIZE):
        chunk = segments[chunk_start: chunk_start + _SEGMENT_CHUNK_SIZE]
        prompt = _build_segment_prompt(chunk, locale, filename)
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 16000,
            "messages": [{"role": "user", "content": prompt}]
        })
        resp = bedrock_client.invoke_model(modelId=AWS_MODEL_ID, body=body)
        raw  = json.loads(resp["body"].read()).get("content", [{}])[0].get("text", "[]")
        all_translations.update(_parse_segment_response(raw))
    return all_translations


def translate_segments_api(client: anthropic.Anthropic, segments: list,
                            locale: str, filename: str) -> dict:
    """Translate structured segments via Anthropic API. Chunks large presentations.
    Returns {id: translated_text}.
    """
    all_translations: dict = {}
    for chunk_start in range(0, len(segments), _SEGMENT_CHUNK_SIZE):
        chunk = segments[chunk_start: chunk_start + _SEGMENT_CHUNK_SIZE]
        prompt = _build_segment_prompt(chunk, locale, filename)
        resp = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=16000,
            messages=[{"role": "user", "content": prompt}]
        )
        all_translations.update(_parse_segment_response(resp.content[0].text))
    return all_translations


def _replace_para_text(para, new_text: str):
    """Replace a paragraph's text while preserving the first run's formatting."""
    runs = para.runs
    if not runs:
        return
    runs[0].text = new_text
    for run in runs[1:]:
        run.text = ""


def rebuild_pptx(original_path: Path, translations: dict, out_path: Path):
    """Write a translated PPTX preserving slide layout, images, tables and run formatting."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(original_path))
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            # ── Regular text frame ──────────────────────────────────────────
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    seg_id = f"s{si}_sh{shi}_p{pi}"
                    if seg_id in translations:
                        _replace_para_text(para, translations[seg_id])
            # ── Table cells ─────────────────────────────────────────────────
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for ri, row in enumerate(shape.table.rows):
                    for ci, cell in enumerate(row.cells):
                        for pi, para in enumerate(cell.text_frame.paragraphs):
                            seg_id = f"s{si}_sh{shi}_r{ri}c{ci}_p{pi}"
                            if seg_id in translations:
                                _replace_para_text(para, translations[seg_id])
    prs.save(str(out_path))


def rebuild_docx(original_path: Path, translations: dict, out_path: Path):
    """Write a translated DOCX preserving paragraph styles, run formatting, and table cells."""
    from docx import Document
    doc = Document(str(original_path))
    # Regular paragraphs
    for pi, para in enumerate(doc.paragraphs):
        seg_id = f"p{pi}"
        if seg_id not in translations or not para.text.strip():
            continue
        new_text = translations[seg_id]
        runs = para.runs
        if not runs:
            continue
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
    # Table cells
    for ti, table in enumerate(doc.tables):
        for ri, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                for pi, para in enumerate(cell.paragraphs):
                    seg_id = f"t{ti}_r{ri}_c{ci}_p{pi}"
                    if seg_id not in translations or not para.text.strip():
                        continue
                    new_text = translations[seg_id]
                    runs = para.runs
                    if not runs:
                        continue
                    runs[0].text = new_text
                    for run in runs[1:]:
                        run.text = ""
    doc.save(str(out_path))


# ── Agent Helper Functions ─────────────────────────────────────────────────────
def _build_agent_user_message(content: str, images: list, locale: str, filename: str) -> list:
    """Build a multimodal message content block list (text + up to 5 images)."""
    lang_info = LANGUAGES.get(locale, {})
    prompt = f"""Translate the following promotional material from English (en-US) to {lang_info.get('name', locale)} ({locale}) for {lang_info.get('country', '')}.

DOCUMENT: {filename}

SOURCE CONTENT:
---
{content[:6000]}
---

Instructions:
1. Translate the full document naturally into the target language.
2. Call flag_segment for regulatory claims, clinical data, safety statements, or uncertain text.
3. Call apply_glossary_term for each brand/product term you preserve.
4. If images are attached, analyze and incorporate their text into the translation.
5. Call submit_translation when the translation is complete."""

    blocks = [{"type": "text", "text": prompt}]
    for img in images[:5]:  # Claude supports up to 5 images per message
        blocks.append({"type": "image",
                        "source": {"type": "base64",
                                   "media_type": img["media_type"],
                                   "data": img["data"]}})
    return blocks


def _process_tool_call(name: str, inp: dict,
                       flagged_segments: list, glossary_terms: list,
                       translation_holder: list) -> str:
    """Execute a tool call, updating state lists in place. Returns result text."""
    if name == "flag_segment":
        flagged_segments.append(inp)
        return "Segment flagged for review."
    elif name == "apply_glossary_term":
        term = inp.get("term", "")
        if term and term not in glossary_terms:
            glossary_terms.append(term)
        return f"Term '{term}' recorded."
    elif name == "submit_translation":
        translation_holder.append(inp)
        return "Translation submitted successfully."
    return "Unknown tool."


def _assemble_agent_result(translation: dict, locale: str,
                           flagged_segments: list, glossary_terms: list) -> dict:
    return {
        "translated_content":     translation.get("translated_content", ""),
        "source_locale":          "en-US",
        "target_locale":          locale,
        "confidence_score":       translation.get("confidence_score", 0.0),
        "flagged_segments":       flagged_segments,
        "glossary_terms_applied": glossary_terms,
        "translation_notes":      translation.get("translation_notes", ""),
    }


def translate_document_agent_bedrock(bedrock_client, content: str, images: list,
                                     locale: str, filename: str) -> dict:
    """Agent-based translation via AWS Bedrock with vision + tool use."""
    messages = [{"role": "user",
                 "content": _build_agent_user_message(content, images, locale, filename)}]
    flagged_segments, glossary_terms, translation_holder = [], [], []

    for _ in range(10):
        body = json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 8000,
            "system": AGENT_SYSTEM_PROMPT,
            "tools": TRANSLATION_TOOLS,
            "messages": messages
        })
        resp      = bedrock_client.invoke_model(modelId=AWS_MODEL_ID, body=body)
        resp_body = json.loads(resp["body"].read())
        blocks    = resp_body.get("content", [])
        messages.append({"role": "assistant", "content": blocks})

        if resp_body.get("stop_reason") != "tool_use":
            break

        tool_results = []
        for block in blocks:
            if block.get("type") != "tool_use":
                continue
            result_text = _process_tool_call(
                block["name"], block["input"],
                flagged_segments, glossary_terms, translation_holder)
            tool_results.append({"type": "tool_result",
                                  "tool_use_id": block["id"],
                                  "content": result_text})
        messages.append({"role": "user", "content": tool_results})
        if translation_holder:
            break

    if not translation_holder:
        raise Exception("Agent did not call submit_translation — check Bedrock model access.")
    return _assemble_agent_result(translation_holder[0], locale, flagged_segments, glossary_terms)


def translate_document_agent(client: anthropic.Anthropic, content: str, images: list,
                              locale: str, filename: str) -> dict:
    """Agent-based translation via Anthropic API with vision + tool use."""
    messages = [{"role": "user",
                 "content": _build_agent_user_message(content, images, locale, filename)}]
    flagged_segments, glossary_terms, translation_holder = [], [], []

    for _ in range(10):
        resp = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=8000,
            system=AGENT_SYSTEM_PROMPT,
            tools=TRANSLATION_TOOLS,
            messages=messages
        )
        # Normalise SDK objects to plain dicts for the message history
        blocks = []
        for b in resp.content:
            if b.type == "text":
                blocks.append({"type": "text", "text": b.text})
            elif b.type == "tool_use":
                blocks.append({"type": "tool_use", "id": b.id,
                                "name": b.name, "input": b.input})
        messages.append({"role": "assistant", "content": blocks})

        if resp.stop_reason != "tool_use":
            break

        tool_results = []
        for block in blocks:
            if block.get("type") != "tool_use":
                continue
            result_text = _process_tool_call(
                block["name"], block["input"],
                flagged_segments, glossary_terms, translation_holder)
            tool_results.append({"type": "tool_result",
                                  "tool_use_id": block["id"],
                                  "content": result_text})
        messages.append({"role": "user", "content": tool_results})
        if translation_holder:
            break

    if not translation_holder:
        raise Exception("Agent did not call submit_translation — check API key and model access.")
    return _assemble_agent_result(translation_holder[0], locale, flagged_segments, glossary_terms)


def scan_source_folder() -> dict:
    """Return {folder_name: [Path, ...]} for all source files."""
    result = {}
    if not SOURCE_DIR.exists():
        SOURCE_DIR.mkdir(parents=True, exist_ok=True)
    for folder in sorted(SOURCE_DIR.iterdir()):
        if folder.is_dir():
            files = sorted([f for f in folder.iterdir()
                            if f.suffix.lower() in (".txt", ".docx", ".pdf", ".pptx", ".ppt")])
            if files:
                result[folder.name] = files
    # Also files in root of source_materials
    root_files = sorted([f for f in SOURCE_DIR.iterdir()
                         if f.is_file() and f.suffix.lower() in (".txt", ".docx", ".pdf", ".pptx", ".ppt")])
    if root_files:
        result["(root)"] = root_files
    return result


def save_translation(file_path: Path, locale: str, content: str, metadata: dict) -> Path:
    """Save translated content to per-locale folder, preserving original file format."""
    locale_info = LANGUAGES.get(locale, {})
    country = locale_info.get("country", locale).replace(" ", "_")
    out_dir = TRANSLATED_DIR / country
    out_dir.mkdir(parents=True, exist_ok=True)
    stem   = file_path.stem
    suffix = file_path.suffix.lower()
    seg_translations = metadata.get("segment_translations")

    # ── Native PPTX output ────────────────────────────────────────────────────
    if suffix in (".pptx", ".ppt") and seg_translations:
        out_path = out_dir / f"{stem}_{locale}.pptx"
        rebuild_pptx(file_path, seg_translations, out_path)
        return out_path

    # ── Native PPTX fallback (no segment translations — basic layout) ─────────
    if suffix in (".pptx", ".ppt"):
        out_path = out_dir / f"{stem}_{locale}.pptx"
        try:
            from pptx import Presentation
            from pptx.util import Pt
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]   # blank
            title_layout = prs.slide_layouts[1]   # title + content
            lines = [l for l in content.split("\n") if l.strip()]
            chunk_size = 20
            for chunk_start in range(0, max(len(lines), 1), chunk_size):
                chunk = lines[chunk_start: chunk_start + chunk_size]
                slide = prs.slides.add_slide(title_layout)
                title_ph = slide.placeholders[0]
                body_ph  = slide.placeholders[1]
                title_ph.text = stem if chunk_start == 0 else f"{stem} (cont.)"
                body_ph.text  = "\n".join(chunk)
            prs.save(str(out_path))
            return out_path
        except Exception:
            pass  # fall through to txt

    # ── Native DOCX output ────────────────────────────────────────────────────
    if suffix == ".docx" and seg_translations:
        out_path = out_dir / f"{stem}_{locale}.docx"
        rebuild_docx(file_path, seg_translations, out_path)
        return out_path

    # ── Native DOCX fallback (no segment translations — plain paragraphs) ─────
    if suffix == ".docx":
        out_path = out_dir / f"{stem}_{locale}.docx"
        try:
            from docx import Document
            doc = Document()
            for para_text in content.split("\n"):
                doc.add_paragraph(para_text if para_text.strip() else "")
            doc.save(str(out_path))
            return out_path
        except Exception:
            pass  # fall through to txt

    # ── Plain-text fallback (txt / pdf / unknown) ─────────────────────────────
    out_path = out_dir / f"{stem}_{locale}.txt"
    full_content = f"""TRANSLATED DOCUMENT
==================
Original: {file_path.name}
Locale: {locale} ({locale_info.get('name', '')} — {locale_info.get('country', '')})
Translated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Model: claude-sonnet-4-5 (AWS Bedrock)
Confidence: {metadata.get('confidence_score', 'N/A')}
Flagged Segments: {len(metadata.get('flagged_segments', []))}
Glossary Terms Applied: {', '.join(metadata.get('glossary_terms_applied', [])) or 'None'}

Translation Notes: {metadata.get('translation_notes', 'N/A')}

{'='*50}
TRANSLATED CONTENT
{'='*50}

{content}
"""
    out_path.write_text(full_content, encoding="utf-8")
    return out_path


def add_audit(action: str, detail: str, level: str = "info"):
    ts = datetime.now().strftime("%H:%M:%S")
    st.session_state.audit_log.insert(0, {
        "time": ts, "action": action, "detail": detail, "level": level
    })
    if len(st.session_state.audit_log) > 100:
        st.session_state.audit_log = st.session_state.audit_log[:100]


def translate_document(client: anthropic.Anthropic, content: str,
                        locale: str, filename: str) -> dict:
    """Call Claude via Anthropic API and return parsed result."""
    lang_info = LANGUAGES.get(locale, {})
    lang_name = lang_info.get("name", locale)
    country = lang_info.get("country", "")

    prompt = f"""Translate the following promotional material from English (en-US) to {lang_name} ({locale}) for {country}.

DOCUMENT: {filename}
TARGET LOCALE: {locale} ({lang_name} — {country})

SOURCE CONTENT:
---
{content[:6000]}
---

Remember: Return ONLY the JSON object as specified. No markdown fencing."""

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4000,
        system=COSTAR_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}]
    )
    raw = response.content[0].text.strip()
    # Strip markdown fences if present
    raw = re.sub(r'^```(?:json)?\s*', '', raw)
    raw = re.sub(r'\s*```$', '', raw)
    result = json.loads(raw)
    return result


def translate_document_bedrock(bedrock_client, content: str,
                                locale: str, filename: str) -> dict:
    """Call Claude via AWS Bedrock and return parsed result."""
    lang_info = LANGUAGES.get(locale, {})
    lang_name = lang_info.get("name", locale)
    country = lang_info.get("country", "")

    prompt = f"""Translate the following promotional material from English (en-US) to {lang_name} ({locale}) for {country}.

DOCUMENT: {filename}
TARGET LOCALE: {locale} ({lang_name} — {country})

SOURCE CONTENT:
---
{content[:6000]}
---

Remember: Return ONLY the JSON object as specified. No markdown fencing."""

    # Prepare the request body for AWS Bedrock
    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 8000,
        "system": COSTAR_SYSTEM_PROMPT,
        "messages": [
            {
                "role": "user",
                "content": prompt
            }
        ]
    })

    try:
        response = bedrock_client.invoke_model(
            modelId=AWS_MODEL_ID,
            body=body
        )
        
        response_body = json.loads(response['body'].read())
        raw = response_body['content'][0]['text'].strip()
        
        # Strip markdown fences if present
        raw = re.sub(r'^```(?:json)?\s*', '', raw)
        raw = re.sub(r'\s*```$', '', raw)
        result = json.loads(raw)
        return result
    except Exception as e:
        raise Exception(f"Error calling AWS Bedrock: {e}")


def build_zip_download() -> bytes:
    """Bundle all translated files into a ZIP."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        if TRANSLATED_DIR.exists():
            for f in TRANSLATED_DIR.rglob("*"):
                if f.is_file():
                    zf.write(f, f.relative_to(TRANSLATED_DIR.parent))
    return buf.getvalue()


def file_type_badge(suffix: str) -> str:
    s = suffix.lstrip(".").lower()
    css = {"txt": "fc-type-txt", "docx": "fc-type-docx",
           "pdf": "fc-type-pdf", "pptx": "fc-type-pptx"}.get(s, "fc-type-txt")
    return f'<span class="fc-type {css}">{s}</span>'


def confidence_color(score: float) -> str:
    if score >= 0.9:  return "#0F6E56"
    if score >= 0.75: return "#BA7517"
    return "#A02D2D"


# ── SIDEBAR ───────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("""
    <div class="logo-block">
      <div class="logo-title">🌐 INT Mkt Translator</div>
      <div class="logo-sub">Translation Agent — POC v1.0</div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("## Configuration")
    
    # LLM Provider Selection
    llm_provider = st.selectbox(
        "LLM Provider",
        options=["AWS Bedrock", "Anthropic API"],
        index=0 if st.session_state.llm_provider == "AWS Bedrock" else 1,
        help="Choose between AWS Bedrock or direct Anthropic API"
    )
    st.session_state.llm_provider = llm_provider
    
    # Initialize AWS Bedrock client if selected
    if llm_provider == "AWS Bedrock":
        if st.session_state.bedrock_client is None:
            try:
                session = boto3.Session(profile_name=AWS_PROFILE_NAME)
                st.session_state.bedrock_client = session.client("bedrock-runtime", region_name=AWS_REGION)
                st.success("✓ AWS Bedrock client initialized")
                add_audit("INIT", f"AWS Bedrock client initialized (profile: {AWS_PROFILE_NAME}, region: {AWS_REGION})")
            except Exception as e:
                st.error(f"✗ Error initializing AWS Bedrock: {e}")
                add_audit("ERROR", f"AWS Bedrock initialization failed: {e}", "error")
    else:
        # Show API Key input for Anthropic API
        api_key = st.text_input(
            "Anthropic API Key",
            value=st.session_state.api_key,
            type="password",
            placeholder="sk-ant-...",
            help="Enter your Anthropic API key for direct API access"
        )
        if api_key:
            st.session_state.api_key = api_key

    st.markdown("---")
    st.markdown("## Agent Mode")
    _mode_opts = ["LangGraph", "Tool-Use Agent", "Direct"]
    agent_mode = st.selectbox(
        "Translation Pipeline",
        options=_mode_opts,
        index=_mode_opts.index(st.session_state.get("agent_mode", "LangGraph")),
        help="LangGraph: 3-node pipeline (Translator→Regulatory→QA)\nTool-Use Agent: single agent with tool calls\nDirect: single-shot JSON (fastest)"
    )
    st.session_state.agent_mode = agent_mode
    _mode_desc = {
        "LangGraph":      "🔗 3-node pipeline · Translator → Regulatory → QA · auto-revision",
        "Tool-Use Agent": "🤖 Single agent · tool calls · vision · iterative flagging",
        "Direct":         "⚡ Single-shot JSON · fastest · text only",
    }
    if not LANGGRAPH_AVAILABLE and agent_mode == "LangGraph":
        st.warning("langgraph not installed. Run: pip install langgraph>=0.2.0")
    st.markdown(f'<p style="font-size:11px;color:#93B4D4 !important">{_mode_desc[agent_mode]}</p>', unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("## Target Languages")
    st.markdown('<p style="font-size:11px;color:#93B4D4 !important">Select locales to translate into</p>', unsafe_allow_html=True)

    default_locales = ["fr-FR", "de-DE", "pt-BR", "ja-JP"]
    selected_locales = []
    cols = st.columns(2)
    for i, (locale, info) in enumerate(LANGUAGES.items()):
        with cols[i % 2]:
            checked = st.checkbox(
                f"{info['flag']} {info['name']}",
                value=locale in default_locales,
                key=f"locale_{locale}"
            )
            if checked:
                selected_locales.append(locale)

    st.markdown("---")
    st.markdown("## Source Folder")
    folder_data = scan_source_folder()
    total_files = sum(len(v) for v in folder_data.values())
    st.markdown(f'<p style="font-size:12px;color:#93B4D4 !important">📁 {len(folder_data)} folders · {total_files} files</p>', unsafe_allow_html=True)

    # Folder tree
    for folder, files in folder_data.items():
        st.markdown(f'<div class="folder-node is-dir">📂 {folder}</div>', unsafe_allow_html=True)
        for f in files:
            key = str(f)
            is_translated = key in st.session_state.translation_results
            style = "is-translated" if is_translated else "is-file"
            icon = "✅" if is_translated else "📄"
            st.markdown(f'<div class="folder-node {style}">{icon} {f.name}</div>', unsafe_allow_html=True)

    st.markdown("---")
    _has_translated = TRANSLATED_DIR.exists() and any(
        f for f in TRANSLATED_DIR.rglob("*") if f.is_file()
    )
    if _has_translated:
        zip_data = build_zip_download()
        st.download_button(
            "⬇️ Download All Translations",
            data=zip_data,
            file_name=f"translations_{datetime.now().strftime('%Y%m%d_%H%M')}.zip",
            mime="application/zip",
            use_container_width=True
        )


# ── MAIN AREA ─────────────────────────────────────────────────────────────────
st.markdown("""
<div class="main-header">
  <h1>INT Marketing Translation Agent</h1>
  <p>POC · AWS Bedrock (Claude Sonnet) · Azure OpenAI Embeddings · AWS OpenSearch · FastAPI · React Native</p>
</div>
""", unsafe_allow_html=True)

# Metrics row
total_docs = sum(len(v) for v in scan_source_folder().values())
done_jobs = len(st.session_state.translation_results)
total_flags = sum(
    len(r.get("flagged_segments", []))
    for results in st.session_state.translation_results.values()
    for r in results.values()
)
total_locales_done = sum(
    len(results) for results in st.session_state.translation_results.values()
)

st.markdown(f"""
<div class="metric-row">
  <div class="metric-card">
    <div class="mc-label">Source Documents</div>
    <div class="mc-val">{total_docs}</div>
    <div class="mc-sub">in source_materials/</div>
  </div>
  <div class="metric-card">
    <div class="mc-label">Files Translated</div>
    <div class="mc-val">{done_jobs}</div>
    <div class="mc-sub">documents processed</div>
  </div>
  <div class="metric-card">
    <div class="mc-label">Locale Copies</div>
    <div class="mc-val">{total_locales_done}</div>
    <div class="mc-sub">translations generated</div>
  </div>
  <div class="metric-card">
    <div class="mc-label">Flagged Segments</div>
    <div class="mc-val" style="color:{'#A02D2D' if total_flags > 0 else '#0F6E56'}">{total_flags}</div>
    <div class="mc-sub">requiring human review</div>
  </div>
</div>
""", unsafe_allow_html=True)

# ── TABS ──────────────────────────────────────────────────────────────────────
tab1, tab2, tab3, tab4 = st.tabs(["📂 Source Files", "🌐 Translate", "📋 Results", "📝 Audit Log"])

# ── TAB 1: SOURCE FILES ───────────────────────────────────────────────────────
with tab1:
    col_left, col_right = st.columns([1, 1.4])

    with col_left:
        st.markdown('<div class="section-title">Source Materials Folder</div>', unsafe_allow_html=True)

        folder_data = scan_source_folder()

        # Upload new files
        with st.expander("➕ Upload New Files", expanded=False):
            uploaded = st.file_uploader(
                "Drop files here",
                accept_multiple_files=True,
                type=["txt", "docx", "pdf", "pptx"],
                label_visibility="collapsed"
            )
            if uploaded:
                folder_choice = st.selectbox(
                    "Save to folder",
                    ["campaign_decks", "product_brochures", "regulatory_docs", "other"]
                )
                if st.button("Save to Source Folder", type="primary"):
                    target_dir = SOURCE_DIR / folder_choice
                    target_dir.mkdir(parents=True, exist_ok=True)
                    for uf in uploaded:
                        dest = target_dir / uf.name
                        content = read_uploaded_content(uf)
                        st.session_state.uploaded_files_cache[str(dest)] = content
                        dest.write_bytes(uf.getvalue())
                        add_audit("UPLOAD", f"{uf.name} → {folder_choice}/")
                    st.success(f"Saved {len(uploaded)} file(s)")
                    st.rerun()

        # List files by folder
        if not folder_data:
            st.info("No source files found. Upload files above or add them to source_materials/")
        else:
            for folder_name, files in folder_data.items():
                st.markdown(f'<div class="section-title">📂 {folder_name}</div>', unsafe_allow_html=True)
                for f in files:
                    is_selected = st.session_state.selected_file == str(f)
                    is_done = str(f) in st.session_state.translation_results
                    done_badge = f'<span class="status-pill status-done">✓ translated</span>' if is_done else ""
                    badge = file_type_badge(f.suffix)
                    card_class = "file-card selected" if is_selected else "file-card"
                    size_kb = f.stat().st_size // 1024 if f.exists() else 0
                    st.markdown(f"""
                    <div class="{card_class}" onclick="">
                      <div class="fc-name">{badge} {f.name} {done_badge}</div>
                      <div class="fc-meta">{size_kb} KB · {f.parent.name}</div>
                    </div>
                    """, unsafe_allow_html=True)
                    if st.button(f"Select → {f.name}", key=f"sel_{f}", use_container_width=True,
                                 type="secondary" if not is_selected else "primary"):
                        st.session_state.selected_file = str(f)
                        st.rerun()

    with col_right:
        st.markdown('<div class="section-title">File Preview</div>', unsafe_allow_html=True)
        if st.session_state.selected_file:
            p = Path(st.session_state.selected_file)
            if p.exists():
                content = read_file_content(p)
                st.markdown(f"""
                <div class="result-card">
                  <div class="rc-header">
                    <span class="rc-locale">{p.name}</span>
                    <span class="status-pill status-pending">en-US source</span>
                  </div>
                </div>
                """, unsafe_allow_html=True)
                st.text_area("Content", value=content, height=500, label_visibility="collapsed")
            else:
                st.warning("File not found on disk.")
        else:
            st.markdown("""
            <div style="text-align:center;padding:4rem 2rem;color:#7B8FA6">
              <div style="font-size:48px;margin-bottom:1rem">📄</div>
              <div style="font-size:15px;font-weight:600">Select a file to preview</div>
              <div style="font-size:13px;margin-top:0.5rem">Click any file on the left to view its contents</div>
            </div>
            """, unsafe_allow_html=True)


# ── TAB 2: TRANSLATE ──────────────────────────────────────────────────────────
with tab2:
    st.markdown('<div class="section-title">Translation Configuration</div>', unsafe_allow_html=True)

    # Provider-specific validation messages
    if st.session_state.llm_provider == "AWS Bedrock":
        if st.session_state.bedrock_client is None:
            st.warning("⚠️ AWS Bedrock client not initialized. Check your AWS configuration.")
    else:
        if not st.session_state.api_key:
            st.warning("⚠️ Enter your Anthropic API key in the sidebar to enable translation.")

    if not selected_locales:
        st.warning("⚠️ Select at least one target language in the sidebar.")

    # Show selected locales
    locale_html = "".join([
        f'<span class="locale-chip active"><span class="flag">{LANGUAGES[l]["flag"]}</span>{LANGUAGES[l]["name"]}</span>'
        for l in selected_locales
    ])
    if locale_html:
        st.markdown(f'<div style="margin-bottom:1rem">{locale_html}</div>', unsafe_allow_html=True)

    col_a, col_b = st.columns([1, 1])

    with col_a:
        st.markdown('<div class="section-title">Select File to Translate</div>', unsafe_allow_html=True)

        folder_data = scan_source_folder()
        all_files = [f for files in folder_data.values() for f in files]
        file_options = {f.name: f for f in all_files}

        if not all_files:
            st.info("No source files. Add files in the Source Files tab first.")
        else:
            chosen_name = st.selectbox("File", list(file_options.keys()), label_visibility="collapsed")
            chosen_file = file_options.get(chosen_name)

            if chosen_file:
                content_preview = read_file_content(chosen_file)
                word_count = len(content_preview.split())
                st.markdown(f"""
                <div class="result-card" style="margin-top:0.5rem">
                  <div class="fc-name">{file_type_badge(chosen_file.suffix)} {chosen_file.name}</div>
                  <div class="fc-meta" style="margin-top:4px">
                    ~{word_count} words · {chosen_file.stat().st_size // 1024} KB · {chosen_file.parent.name}
                  </div>
                </div>
                """, unsafe_allow_html=True)

    with col_b:
        st.markdown('<div class="section-title">Translation Options</div>', unsafe_allow_html=True)
        translate_mode = st.radio(
            "Mode",
            ["Translate selected file only", "Translate ALL source files"],
            horizontal=True,
            label_visibility="collapsed"
        )
        show_raw_json = st.checkbox("Show raw JSON response", value=False)

    st.markdown("---")

    # Translate button
    # Validation depends on provider
    if st.session_state.llm_provider == "AWS Bedrock":
        can_translate = (
            st.session_state.bedrock_client is not None and
            selected_locales and
            all_files and
            not st.session_state.job_running
        )
    else:  # Anthropic API
        can_translate = (
            st.session_state.api_key and
            selected_locales and
            all_files and
            not st.session_state.job_running
        )

    btn_label = (
        f"🚀 Translate into {len(selected_locales)} Language{'s' if len(selected_locales) > 1 else ''}"
        if translate_mode == "Translate selected file only"
        else f"🚀 Translate ALL {len(all_files)} Files × {len(selected_locales)} Languages"
    )

    if st.button(btn_label, type="primary", disabled=not can_translate, use_container_width=True):
        st.session_state.job_running = True

        files_to_process = (
            [chosen_file] if translate_mode == "Translate selected file only"
            else all_files
        )

        # Initialize client based on provider
        if st.session_state.llm_provider == "AWS Bedrock":
            client = st.session_state.bedrock_client
            provider_name = "AWS Bedrock"
        else:
            client = anthropic.Anthropic(api_key=st.session_state.api_key)
            provider_name = "Anthropic API"
        
        add_audit("START", f"Translation job started using {provider_name} [{st.session_state.agent_mode}]")
        
        progress_bar = st.progress(0)
        status_text = st.empty()
        total_ops = len(files_to_process) * len(selected_locales)
        op_count = 0
        errors = []

        for file_path in files_to_process:
            file_key = str(file_path)
            if file_key not in st.session_state.translation_results:
                st.session_state.translation_results[file_key] = {}

            file_content = read_file_content(file_path)
            add_audit("START", f"Processing {file_path.name} → {len(selected_locales)} locales")

            for locale in selected_locales:
                lang_info = LANGUAGES[locale]
                status_text.markdown(f"""
                <div style="font-size:14px;color:#185FA5;padding:0.5rem 0">
                  ⏳ Translating <strong>{file_path.name}</strong> → {lang_info['flag']} {lang_info['name']} ({locale})
                </div>
                """, unsafe_allow_html=True)

                try:
                    _mode = st.session_state.agent_mode
                    if _mode == "LangGraph":
                        _, images = read_file_with_images(file_path)
                        if st.session_state.llm_provider == "AWS Bedrock":
                            result = translate_document_langgraph_bedrock(
                                client, file_content, images, locale, file_path.name)
                        else:
                            result = translate_document_langgraph(
                                client, file_content, images, locale, file_path.name)
                    elif _mode == "Tool-Use Agent":
                        _, images = read_file_with_images(file_path)
                        if st.session_state.llm_provider == "AWS Bedrock":
                            result = translate_document_agent_bedrock(
                                client, file_content, images, locale, file_path.name)
                        else:
                            result = translate_document_agent(
                                client, file_content, images, locale, file_path.name)
                    else:  # Direct
                        if st.session_state.llm_provider == "AWS Bedrock":
                            result = translate_document_bedrock(client, file_content, locale, file_path.name)
                        else:
                            result = translate_document(client, file_content, locale, file_path.name)

                    # ── Enrich with per-segment translations for native format output ──
                    _suffix = file_path.suffix.lower()
                    if _suffix in ('.pptx', '.docx'):
                        try:
                            _segs = (extract_pptx_segments(file_path)
                                     if _suffix == '.pptx'
                                     else extract_docx_segments(file_path))
                            if _segs:
                                if st.session_state.llm_provider == "AWS Bedrock":
                                    result["segment_translations"] = translate_segments_bedrock(
                                        client, _segs, locale, file_path.name)
                                else:
                                    result["segment_translations"] = translate_segments_api(
                                        client, _segs, locale, file_path.name)
                        except Exception as _seg_err:
                            # Non-fatal: fall back to txt output; log so user can see why
                            result.setdefault("segment_translations", None)
                            add_audit("WARN",
                                      f"Native format rebuild skipped for {file_path.name} "
                                      f"({type(_seg_err).__name__}: {_seg_err})",
                                      "warning")

                    st.session_state.translation_results[file_key][locale] = result

                    # Save to disk
                    translated_text = result.get("translated_content", "")
                    saved_path = save_translation(file_path, locale, translated_text, result)
                    result["saved_path"] = str(saved_path)  # keep reference for download

                    n_flags = len(result.get("flagged_segments", []))
                    conf = result.get("confidence_score", 0)
                    add_audit(
                        "DONE",
                        f"{file_path.name} → {locale}: conf={conf:.2f}, flags={n_flags}, saved to {saved_path.parent.name}/",
                        "warning" if n_flags > 0 else "info"
                    )
                except json.JSONDecodeError as e:
                    errors.append(f"{file_path.name}/{locale}: JSON parse error — {e}")
                    add_audit("ERROR", f"{file_path.name}/{locale}: JSON parse failed", "error")
                except Exception as e:
                    errors.append(f"{file_path.name}/{locale}: {e}")
                    add_audit("ERROR", f"{file_path.name}/{locale}: {e}", "error")

                op_count += 1
                progress_bar.progress(op_count / total_ops)

        st.session_state.job_running = False
        status_text.empty()
        progress_bar.empty()

        if errors:
            st.error("Some translations failed:\n" + "\n".join(errors))
        else:
            st.success(f"✅ Translation complete! {op_count} locale copies generated across {len(files_to_process)} file(s).")
            time.sleep(0.5)
            st.rerun()


# ── TAB 3: RESULTS ────────────────────────────────────────────────────────────
with tab3:
    if not st.session_state.translation_results:
        st.markdown("""
        <div style="text-align:center;padding:4rem 2rem;color:#7B8FA6">
          <div style="font-size:48px;margin-bottom:1rem">🌐</div>
          <div style="font-size:15px;font-weight:600">No translations yet</div>
          <div style="font-size:13px;margin-top:0.5rem">Go to the Translate tab and run a translation</div>
        </div>
        """, unsafe_allow_html=True)
    else:
        # File selector
        result_files = list(st.session_state.translation_results.keys())
        result_file_names = {Path(f).name: f for f in result_files}
        chosen_result = st.selectbox("View results for:", list(result_file_names.keys()))
        chosen_key = result_file_names.get(chosen_result)

        if chosen_key:
            results_for_file = st.session_state.translation_results[chosen_key]
            st.markdown(f'<div class="section-title">{len(results_for_file)} translations · {Path(chosen_key).name}</div>', unsafe_allow_html=True)

            # Summary row
            all_flags = [f for r in results_for_file.values() for f in r.get("flagged_segments", [])]
            avg_conf = sum(r.get("confidence_score", 0) for r in results_for_file.values()) / max(len(results_for_file), 1)
            st.markdown(f"""
            <div class="metric-row" style="margin-bottom:1rem">
              <div class="metric-card">
                <div class="mc-label">Locales Generated</div>
                <div class="mc-val">{len(results_for_file)}</div>
              </div>
              <div class="metric-card">
                <div class="mc-label">Avg Confidence</div>
                <div class="mc-val" style="color:{confidence_color(avg_conf)}">{avg_conf:.0%}</div>
              </div>
              <div class="metric-card">
                <div class="mc-label">Total Flagged</div>
                <div class="mc-val" style="color:{'#A02D2D' if all_flags else '#0F6E56'}">{len(all_flags)}</div>
                <div class="mc-sub">require human review</div>
              </div>
              <div class="metric-card">
                <div class="mc-label">Glossary Terms</div>
                <div class="mc-val">{len(set(t for r in results_for_file.values() for t in r.get('glossary_terms_applied', [])))}</div>
                <div class="mc-sub">brand terms preserved</div>
              </div>
            </div>
            """, unsafe_allow_html=True)

            # Per-locale results
            for locale, result in results_for_file.items():
                lang_info = LANGUAGES.get(locale, {})
                flags = result.get("flagged_segments", [])
                conf = result.get("confidence_score", 0)
                flag_badge = (f'<span class="status-pill status-flagged">⚠️ {len(flags)} flagged</span>'
                              if flags else '<span class="status-pill status-done">✓ clean</span>')
                conf_color = confidence_color(conf)

                with st.expander(
                    f"{lang_info.get('flag','🌐')} {lang_info.get('name', locale)} ({locale}) — {conf:.0%} confidence",
                    expanded=(locale in ["fr-FR", "de-DE"])
                ):
                    col1, col2 = st.columns([3, 1])
                    with col1:
                        st.markdown(f"""
                        <div class="rc-header">
                          <span class="rc-locale">{lang_info.get('flag','')} {lang_info.get('name','')} — {lang_info.get('country','')}</span>
                          {flag_badge}
                        </div>
                        """, unsafe_allow_html=True)
                        translated_text = result.get("translated_content", "No content")
                        st.text_area(
                            f"Translation ({locale})",
                            value=translated_text,
                            height=300,
                            key=f"ta_{chosen_key}_{locale}",
                            label_visibility="collapsed"
                        )

                        # Download this translation
                        saved_path_str = result.get("saved_path", "")
                        saved_path_obj = Path(saved_path_str) if saved_path_str else None
                        if saved_path_obj and saved_path_obj.exists() and saved_path_obj.suffix.lower() != ".txt":
                            # Serve the native file (PPTX / DOCX) directly
                            _mime_map = {
                                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            }
                            _mime = _mime_map.get(saved_path_obj.suffix.lower(), "application/octet-stream")
                            st.download_button(
                                f"⬇️ Download {locale} translation ({saved_path_obj.suffix[1:].upper()})",
                                data=saved_path_obj.read_bytes(),
                                file_name=saved_path_obj.name,
                                mime=_mime,
                                key=f"dl_{chosen_key}_{locale}"
                            )
                        else:
                            st.download_button(
                                f"⬇️ Download {locale} translation",
                                data=translated_text.encode("utf-8"),
                                file_name=f"{Path(chosen_key).stem}_{locale}.txt",
                                mime="text/plain",
                                key=f"dl_{chosen_key}_{locale}"
                            )

                    with col2:
                        st.markdown(f"""
                        <div class="result-card" style="margin-bottom:0.5rem">
                          <div class="mc-label">Confidence</div>
                          <div style="font-size:22px;font-weight:700;color:{conf_color}">{conf:.0%}</div>
                          <div class="conf-bar-wrap"><div class="conf-bar" style="width:{conf*100:.0f}%;background:{conf_color}"></div></div>
                        </div>
                        """, unsafe_allow_html=True)

                        glossary_terms = result.get("glossary_terms_applied", [])
                        if glossary_terms:
                            st.markdown('<div class="mc-label" style="margin-top:0.5rem">Glossary Applied</div>', unsafe_allow_html=True)
                            for term in glossary_terms[:8]:
                                st.markdown(f'<div style="font-size:12px;color:#185FA5;padding:2px 0">• {term}</div>', unsafe_allow_html=True)

                        notes = result.get("translation_notes", "")
                        if notes:
                            st.markdown('<div class="mc-label" style="margin-top:0.75rem">Reviewer Notes</div>', unsafe_allow_html=True)
                            st.markdown(f'<div style="font-size:12px;color:#4A5568;line-height:1.5">{notes}</div>', unsafe_allow_html=True)

                    # Flagged segments
                    if flags:
                        st.markdown('<div class="section-title" style="margin-top:0.75rem">⚠️ Flagged Segments — Human Review Required</div>', unsafe_allow_html=True)
                        for seg in flags:
                            reason_colors = {
                                "regulatory_claim": "#F8D7DA",
                                "clinical_data": "#FFF3CD",
                                "low_confidence": "#D1ECF1",
                                "brand_term_missing": "#FFE8D6"
                            }
                            bg = reason_colors.get(seg.get("reason", ""), "#F8D7DA")
                            st.markdown(f"""
                            <div class="flag-item" style="background:{bg}">
                              <span class="fi-reason">[{seg.get('reason','?').upper().replace('_',' ')}]</span>
                              <strong>Original:</strong> {seg.get('original', '')[:120]}<br>
                              <strong>Translation:</strong> {seg.get('translation', '')[:120]}<br>
                              <em style="font-size:11px;color:#666">Confidence: {seg.get('confidence', 0):.0%}</em>
                            </div>
                            """, unsafe_allow_html=True)

                    # Raw JSON toggle
                    if show_raw_json:
                        with st.expander("Raw JSON response"):
                            st.json(result)


# ── TAB 4: AUDIT LOG ─────────────────────────────────────────────────────────
with tab4:
    col_al, col_ar = st.columns([3, 1])
    with col_al:
        st.markdown('<div class="section-title">Translation Audit Trail</div>', unsafe_allow_html=True)
    with col_ar:
        if st.button("🗑 Clear Log", use_container_width=True):
            st.session_state.audit_log = []
            st.rerun()

    if not st.session_state.audit_log:
        st.markdown('<div style="color:#7B8FA6;font-size:14px;padding:2rem 0">No activity yet. Audit entries appear here when translations are run.</div>', unsafe_allow_html=True)
    else:
        level_colors = {"info": "#185FA5", "warning": "#BA7517", "error": "#A02D2D"}
        for entry in st.session_state.audit_log:
            color = level_colors.get(entry.get("level", "info"), "#185FA5")
            st.markdown(f"""
            <div class="audit-row">
              <span class="ar-time">{entry['time']}</span>
              <span class="ar-action" style="color:{color}">{entry['action']}</span>
              <span>{entry['detail']}</span>
            </div>
            """, unsafe_allow_html=True)

    # Translated folder structure
    st.markdown('<div class="section-title" style="margin-top:1.5rem">Translated Output Folder</div>', unsafe_allow_html=True)
    if TRANSLATED_DIR.exists():
        _ext_icons = {".txt": "📄", ".docx": "📝", ".pptx": "📊", ".pdf": "📕"}
        for country_dir in sorted(TRANSLATED_DIR.iterdir()):
            if country_dir.is_dir():
                files = [f for f in country_dir.iterdir() if f.is_file()]
                st.markdown(f'<div class="folder-node is-dir">📂 translated/{country_dir.name}/ ({len(files)} files)</div>', unsafe_allow_html=True)
                for tf in sorted(files):
                    size_kb = tf.stat().st_size // 1024
                    icon = _ext_icons.get(tf.suffix.lower(), "📄")
                    st.markdown(f'<div class="folder-node is-translated">{icon} {tf.name} ({size_kb} KB)</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div style="color:#7B8FA6;font-size:13px">No translated files yet.</div>', unsafe_allow_html=True)
